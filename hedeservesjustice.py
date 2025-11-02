# generate_hedeservesjustice_templates_fixed.py
# Fixed: use prs.slide_width/slide_height; pass prs into helper functions.

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_LINE
from pptx.enum.shapes import MSO_SHAPE

EMU_PER_INCH = 914400

def set_slide_size(prs, w_in, h_in):
    prs.slide_width = int(w_in * EMU_PER_INCH)
    prs.slide_height = int(h_in * EMU_PER_INCH)

def add_textbox(slide, left_in, top_in, width_in, height_in, text, font_size=32, bold=False, color=(0,0,0), align="left"):
    box = slide.shapes.add_textbox(int(left_in*EMU_PER_INCH), int(top_in*EMU_PER_INCH),
                                   int(width_in*EMU_PER_INCH), int(height_in*EMU_PER_INCH))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    font = run.font
    font.size = Pt(font_size)
    font.bold = bold
    font.color.rgb = RGBColor(*color)
    p.alignment = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}[align]
    return box

def add_safe_zone_guides(slide, prs, mode, margin_ratio=0.0556):
    # Use presentation-wide size in EMU
    w = prs.slide_width
    h = prs.slide_height
    mx = int(w * margin_ratio)
    my = int(h * margin_ratio)
    rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, mx, my, w - 2*mx, h - 2*my)
    rect.fill.background()
    rect.line.color.rgb = RGBColor(200,200,200)
    rect.line.width = Pt(2)
    rect.line.dash_style = MSO_LINE.DASH

    if mode == "9x16":
        reels_h = int(h * 0.75)
        reels_top = int((h - reels_h) / 2)
        r2 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, mx, reels_top, w - 2*mx, reels_h)
        r2.fill.background()
        r2.line.color.rgb = RGBColor(255,0,0)
        r2.line.width = Pt(2)
        r2.line.dash_style = MSO_LINE.DASH

        stories_h = int(h * 0.8385)
        stories_top = int((h - stories_h) / 2)
        r3 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, mx, stories_top, w - 2*mx, stories_h)
        r3.fill.background()
        r3.line.color.rgb = RGBColor(0,128,255)
        r3.line.width = Pt(2)
        r3.line.dash_style = MSO_LINE.DASH

def add_brand_footer(slide, prs, handle="@hedeservesjustice"):
    w_in = prs.slide_width / EMU_PER_INCH
    h_in = prs.slide_height / EMU_PER_INCH
    add_textbox(slide, w_in - 3.9, h_in - 0.8, 3.5, 0.5, handle, font_size=16, bold=True, color=(30,30,30), align="right")

def add_title_punch(slide, prs, mode):
    add_safe_zone_guides(slide, prs, mode)
    add_textbox(slide, 0.7, 1.0, 7.5, 1.2, "Headline (≤ 8 words)", font_size=56, bold=True)
    add_textbox(slide, 0.7, 2.4, 7.5, 1.0, "Subhead (1 line)", font_size=28)
    add_textbox(slide, 0.7, 3.6, 7.5, 2.5, "Body: 2–4 short bullets, large mobile-friendly text.", font_size=24)
    add_brand_footer(slide, prs)

def add_myth_fact(slide, prs, mode):
    add_safe_zone_guides(slide, prs, mode)
    add_textbox(slide, 0.7, 0.9, 7.5, 0.8, "Myth:", font_size=36, bold=True, color=(180,0,0))
    add_textbox(slide, 0.7, 1.6, 7.5, 2.4, "Write the misconception in plain language.")
    add_textbox(slide, 0.7, 4.3, 7.5, 0.8, "Fact:", font_size=36, bold=True, color=(0,120,0))
    add_textbox(slide, 0.7, 5.0, 7.5, 2.4, "Write the verified fact and cite law or source.")
    add_brand_footer(slide, prs)

def add_rights_remedies(slide, prs, mode):
    add_safe_zone_guides(slide, prs, mode)
    add_textbox(slide, 0.7, 0.9, 7.5, 1.0, "Rights & Remedies", font_size=40, bold=True)
    add_textbox(slide, 0.9, 2.0, 7.2, 3.8, "1) Step one (who/where)\n2) Step two (forms/fees)\n3) Step three (timelines)\n4) Step four (evidence)", font_size=28)
    add_brand_footer(slide, prs)

def add_case_snapshot(slide, prs, mode):
    add_safe_zone_guides(slide, prs, mode)
    add_textbox(slide, 0.7, 0.9, 7.5, 0.8, "Case snapshot", font_size=36, bold=True)
    add_textbox(slide, 0.7, 1.9, 7.5, 1.4, "Issue:", font_size=30, bold=True)
    add_textbox(slide, 0.7, 3.1, 7.5, 1.4, "What the law says:", font_size=30, bold=True)
    add_textbox(slide, 0.7, 4.3, 7.5, 1.4, "Action to take:", font_size=30, bold=True)
    add_brand_footer(slide, prs)

def add_cta(slide, prs, mode):
    add_safe_zone_guides(slide, prs, mode)
    add_textbox(slide, 0.7, 2.2, 7.5, 1.2, "Save • Share • Follow", font_size=48, bold=True, align="center")
    add_textbox(slide, 1.2, 3.6, 6.5, 0.8, "Read more in caption • Not legal advice", font_size=24, align="center")
    add_brand_footer(slide, prs)

def build(filename, mode):
    prs = Presentation()
    if mode == "4x5":
        set_slide_size(prs, 8, 10)    # 4:5 portrait
    elif mode == "1x1":
        set_slide_size(prs, 10, 10)   # square
    elif mode == "9x16":
        set_slide_size(prs, 9, 16)    # story/reel
    elif mode == "1x1.55":
        set_slide_size(prs, 4.2, 6.51)  # ~1:1.55

    blank = prs.slide_layouts[6]
    for fn in (add_title_punch, add_myth_fact, add_rights_remedies, add_case_snapshot, add_cta):
        slide = prs.slides.add_slide(blank)
        fn(slide, prs, mode)

    prs.save(filename)

build('hedeservesjustice_feed_4x5_template.pptx', '4x5')
build('hedeservesjustice_square_1x1_template.pptx', '1x1')
build('hedeservesjustice_story_reel_9x16_template.pptx', '9x16')
build('hedeservesjustice_reel_cover_1x1_55_template.pptx', '1x1.55')
